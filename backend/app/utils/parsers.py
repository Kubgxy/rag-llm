import os
import re
import sqlite3
from typing import Dict, Optional
import pandas as pd

def extract_text_from_docx(file_path: str) -> Dict[int, str]:
    """
    สกัดข้อความจากไฟล์ Word (.docx) โดยย่อหน้าและตารางมาจัดกลุ่มเป็นหน้าจำลอง (Pseudo-pages)
    คืนค่าเป็น Dictionary {เลขหน้าจำลอง: ข้อความ}
    """
    import docx
    pages_text = {}
    doc = docx.Document(file_path)
    
    current_page = 1
    current_text_block = []
    current_char_count = 0
    MAX_CHAR_PER_PAGE = 1200  # เป้าหมายตัวอักษรต่อ 1 หน้าจำลอง (1,000 - 1,500 ตัวอักษร)

    # 1. วนลูปอ่านทุกองค์ประกอบใน body ของเอกสาร (Paragraphs และ Tables)
    # เพื่อรักษาระเบียบการวางลำดับเนื้อหาในเอกสาร
    for element in doc.element.body:
        # ตรวจว่าเป็นย่อหน้า (Paragraph)
        if element.tag.endswith('p'):
            p_id = element.attrib.get('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}paragraphId')
            # ค้นหาย่อหน้าที่มี id ตรงกันในเอกสาร
            for paragraph in doc.paragraphs:
                # ป้องกันกรณีหา paragraphId ไม่เจอหรือเข้าถึงไม่ได้ ให้เปรียบเทียบข้อความหรือใช้ออบเจ็กต์โดยตรง
                if paragraph.text.strip():
                    # สกัดเฉพาะย่อหน้าที่ตรงกัน (หรือถ้าหาแบบละเอียดไม่ได้ ให้รันวนปกติ)
                    pass
            
    # วิธีรันแบบเรียบง่ายและเสถียรที่สุดสำหรับ python-docx คืออ่านแยก Paragraphs และ Tables
    # แต่เนื่องจากเราต้องการแบ่งหน้าสลับตามเนื้อหาจริง เราจะอ่าน Paragraphs ก่อน
    # แล้วสอดแทรกตาราง หรืออ่านวนย่อยดังนี้:
    
    # วนอ่านข้อความทุก Paragraph
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
            
        text_len = len(text)
        # ถ้าข้อความใหม่บวกกับของเก่าแล้วเกินขีดจำกัดหน้า ให้ปัดขึ้นหน้าใหม่
        if current_char_count > 0 and (current_char_count + text_len > MAX_CHAR_PER_PAGE):
            pages_text[current_page] = "\n".join(current_text_block)
            current_page += 1
            current_text_block = []
            current_char_count = 0
            
        current_text_block.append(text)
        current_char_count += text_len
        
    # วนอ่านข้อความในตารางทั้งหมด และต่อข้อความเป็นหน้าจำลองท้ายสุด
    for table in doc.tables:
        table_rows = []
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                table_rows.append(" | ".join(row_text))
                
        if table_rows:
            table_text = "\n".join(table_rows)
            text_len = len(table_text)
            
            if current_char_count > 0 and (current_char_count + text_len > MAX_CHAR_PER_PAGE):
                pages_text[current_page] = "\n".join(current_text_block)
                current_page += 1
                current_text_block = []
                current_char_count = 0
                
            current_text_block.append("[ตารางข้อมูล]\n" + table_text)
            current_char_count += text_len + 12

    # เก็บหน้าสุดท้ายที่เหลืออยู่
    if current_text_block:
        pages_text[current_page] = "\n".join(current_text_block)
        
    # ป้องกันคืนค่าดิบว่างเปล่า
    if not pages_text:
        pages_text[1] = "เอกสารนี้ไม่มีเนื้อหาข้อความ"
        
    return pages_text


def extract_text_from_pptx(file_path: str) -> Dict[int, str]:
    """
    สกัดข้อความจากสไลด์ PowerPoint (.pptx)
    โดยกำหนดให้ 1 Slide = 1 Page จำลอง เพื่อให้การอ้างอิงชัดเจน
    """
    from pptx import Presentation
    pages_text = {}
    prs = Presentation(file_path)
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_contents = []
        
        # ค้นหาข้อความในทุกๆ shapes และตารางที่อยู่ในสไลด์
        for shape in slide.shapes:
            # 1. ข้อความทั่วไปในกล่องข้อความ
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_contents.append(text)
            
            # 2. ข้อความภายในตารางบนสไลด์
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        slide_contents.append(" | ".join(row_text))
                        
        # บันทึกข้อมูลข้อความในสไลด์นั้นๆ
        combined_text = "\n".join(slide_contents)
        if combined_text.strip():
            pages_text[slide_num] = combined_text
        else:
            pages_text[slide_num] = f"[สไลด์หน้า {slide_num} ไม่มีข้อความ หรือเป็นภาพประกอบ]"
            
    if not pages_text:
        pages_text[1] = "ไม่มีข้อมูลสไลด์ในไฟล์นี้"
        
    return pages_text


def extract_text_from_txt_md(file_path: str) -> Dict[int, str]:
    """
    สกัดข้อความจากไฟล์ Plain Text (.txt) หรือ Markdown (.md)
    โดยแบ่งกลุ่มตาม Paragraphs ทุกๆ 1,000 - 1,500 ตัวอักษร
    """
    pages_text = {}
    
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
    except UnicodeDecodeError:
        # Fallback ในกรณีที่ไฟล์ไม่ได้เป็น UTF-8 (เช่น tis-620 หรือ windows-874 ในไทย)
        with open(file_path, "r", encoding="windows-874", errors="ignore") as f:
            content = f.read()
            
    # แยกย่อหน้าด้วย double newline
    paragraphs = [p.strip() for p in re.split(r'\n\s*\n', content) if p.strip()]
    
    current_page = 1
    current_text_block = []
    current_char_count = 0
    MAX_CHAR_PER_PAGE = 1200

    for para in paragraphs:
        text_len = len(para)
        
        # ตรวจจับหัวข้อ Markdown (# ## ###) เพื่อเป็นสัญญาณในการขึ้นหน้าจำลองใหม่หากหน้านั้นมีข้อความพอประมาณแล้ว
        is_header = para.startswith('#')
        
        if current_char_count > 0 and ((current_char_count + text_len > MAX_CHAR_PER_PAGE) or (is_header and current_char_count > 400)):
            pages_text[current_page] = "\n\n".join(current_text_block)
            current_page += 1
            current_text_block = []
            current_char_count = 0
            
        current_text_block.append(para)
        current_char_count += text_len
        
    if current_text_block:
        pages_text[current_page] = "\n\n".join(current_text_block)
        
    if not pages_text:
        pages_text[1] = "ไฟล์ว่างเปล่า ไม่มีข้อความข้อมูล"
        
    return pages_text


def ingest_tabular_file(file_path: str, filename: str, session_id: str) -> int:
    """
    ประมวลผลไฟล์ CSV หรือ Excel (.xlsx) เพื่อบันทึกลงใน SQLite Database
    แยกชื่อตารางตาม Session ID และชื่อไฟล์เพื่อป้องกันการชนกันของข้อมูล
    คืนค่าเป็นจำนวนแถวที่บันทึกสำเร็จ
    """
    # 1. โหลดข้อมูลลง Pandas DataFrame
    ext = os.path.splitext(filename)[1].lower()
    if ext == '.csv':
        try:
            df = pd.read_csv(file_path, encoding='utf-8')
        except UnicodeDecodeError:
            df = pd.read_csv(file_path, encoding='tis-620')
    elif ext in ['.xlsx', '.xls']:
        df = pd.read_excel(file_path)
    else:
        raise ValueError(f"ชนิดไฟล์ตารางไม่ถูกต้อง: {ext}")
        
    # 2. ปรับแต่งชื่อคอลัมน์ให้อยู่ในรูปแบบที่ถูกต้องสำหรับ SQLite (ตัวพิมพ์เล็ก, ไม่มีเว้นวรรค, ไม่มีอักขระพิเศษ)
    clean_cols = []
    for col in df.columns:
        col_str = str(col).strip().lower()
        col_str = re.sub(r'[^a-z0-9_ก-๙]', '_', col_str)  # รองรับอักษรไทยในคอลัมน์เบื้องต้น
        col_str = re.sub(r'_+', '_', col_str)
        col_str = col_str.strip('_')
        if not col_str:
            col_str = f"col_{len(clean_cols)}"
        clean_cols.append(col_str)
    df.columns = clean_cols

    # 3. สร้างชื่อตารางที่ปลอดภัย
    # แปลงอักษรพิเศษในชื่อไฟล์ออก
    safe_filename = re.sub(r'[^a-zA-Z0-9_]', '_', os.path.splitext(filename)[0]).lower()
    safe_session = re.sub(r'[^a-zA-Z0-9_]', '_', session_id).lower()
    table_name = f"table_{safe_session}_{safe_filename}"
    
    # 4. บันทึกลง SQLite
    os.makedirs("data", exist_ok=True)
    db_path = os.path.join("data", "tabular.db")
    
    conn = sqlite3.connect(db_path)
    try:
        # บันทึกข้อมูลลงฐานข้อมูล (ถ้ามีอยู่แล้วให้แทนที่)
        df.to_sql(table_name, conn, if_exists="replace", index=False)
        print(f"📊 [SQLite Ingestion] บันทึกตารางสำเร็จ: {table_name} ({len(df)} แถว, {len(df.columns)} คอลัมน์)")
    finally:
        conn.close()
        
    return len(df)
